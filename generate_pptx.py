"""
Generate PPTX with screenshots of MAVIS 2.0 prototype screens.
Uses Playwright for screenshotting and python-pptx for slide creation.
"""

import os
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "screenshots"
PPTX_PATH = BASE_DIR.parent / "MAVIS 2.0 - Platform Preview.pptx"

# Screen definitions: (filename, title, subtitle, viewport_width, viewport_height)
SCREENS = [
    ("01-dashboard.html", "Dashboard & AI Insights", "Real-time overview with AI-powered pattern detection", 1440, 1200),
    ("02-report-list.html", "Report Management", "Search, filter, and manage all inspection reports", 1440, 900),
    ("03-create-report-details.html", "Create Report — Details", "Step 1: Select report type, equipment auto-populates", 1440, 1100),
    ("04-create-report-findings.html", "Create Report — Findings", "Step 2: Photos, annotations, findings & recommendations", 1440, 1300),
    ("05-create-report-review.html", "Create Report — Review", "Step 3: Review everything before submission", 1440, 1200),
    ("06-mobile-view.html", "Mobile & Offline", "Full report creation from phone — works offline", 500, 960),
]

SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)


def take_screenshots():
    """Take screenshots of all HTML screens using Playwright."""
    OUTPUT_DIR.mkdir(exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch()

        for filename, title, _, vw, vh in SCREENS:
            filepath = BASE_DIR / filename
            if not filepath.exists():
                print(f"  Skipping {filename} — not found")
                continue

            page = browser.new_page(viewport={"width": vw, "height": vh})
            page.goto(f"file://{filepath.resolve()}")
            page.wait_for_load_state("networkidle")
            # Wait for fonts to load
            page.wait_for_timeout(1500)

            screenshot_path = OUTPUT_DIR / f"{Path(filename).stem}.png"

            # For mobile view, screenshot just the device frame
            if "mobile" in filename:
                page.screenshot(path=str(screenshot_path), full_page=False)
            else:
                page.screenshot(path=str(screenshot_path), full_page=True)

            print(f"  Captured: {screenshot_path.name} ({vw}x{vh})")
            page.close()

        browser.close()


def create_pptx():
    """Create a PPTX with the screenshots."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ---- Title Slide ----
    slide = prs.slides.add_slide(blank_layout)

    # Background - dark
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1D, 0x1D, 0x1F)

    # MAVIS 2.0 title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "MAVIS 2.0"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Inter"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "Inspection Reporting Platform — UI Preview"
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0x86, 0x86, 0x8B)
    run2.font.name = "Inter"

    # Accent line
    from pptx.util import Pt as PtSize
    line_shape = slide.shapes.add_shape(
        1, Inches(1.5), Inches(3.3), Inches(1.2), Pt(3)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x94, 0x88)
    line_shape.line.fill.background()

    # "Prepared by Website Artisan" at bottom
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(5), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Prepared by Website Artisan  |  April 2026"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x64, 0x64, 0x68)
    run3.font.name = "Inter"

    # ---- Screenshot Slides ----
    for filename, title, subtitle, vw, vh in SCREENS:
        screenshot_path = OUTPUT_DIR / f"{Path(filename).stem}.png"
        if not screenshot_path.exists():
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Light background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)

        # Title text - top left
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1D, 0x1D, 0x1F)
        run.font.name = "Inter"

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(8), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0x86, 0x86, 0x8B)
        run2.font.name = "Inter"

        # Screenshot image — centered, with shadow effect via positioning
        from PIL import Image
        img = Image.open(screenshot_path)
        img_w, img_h = img.size
        img.close()

        is_mobile = "mobile" in filename

        if is_mobile:
            # Mobile: show centered, smaller
            max_h = Inches(6.0)
            max_w = Inches(3.2)
            ratio = min(max_w / Emu(img_w * 914400 // 96), max_h / Emu(img_h * 914400 // 96))
            pic_w = int(Emu(img_w * 914400 // 96) * ratio)
            pic_h = int(Emu(img_h * 914400 // 96) * ratio)
            pic_left = int((SLIDE_WIDTH - pic_w) / 2)
            pic_top = Inches(1.3)
        else:
            # Desktop: fill width with padding, add rounded appearance
            margin = Inches(0.6)
            available_w = SLIDE_WIDTH - 2 * margin
            available_h = SLIDE_HEIGHT - Inches(1.5)

            scale_w = available_w / Emu(img_w * 914400 // 96)
            scale_h = available_h / Emu(img_h * 914400 // 96)
            ratio = min(scale_w, scale_h)

            pic_w = int(Emu(img_w * 914400 // 96) * ratio)
            pic_h = int(Emu(img_h * 914400 // 96) * ratio)
            pic_left = int((SLIDE_WIDTH - pic_w) / 2)
            pic_top = Inches(1.25)

        pic = slide.shapes.add_picture(str(screenshot_path), pic_left, pic_top, pic_w, pic_h)

        # Add a subtle shadow to the picture
        from pptx.oxml.ns import qn
        spPr = pic._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = pic._element.find(qn('pic:spPr'))

    # ---- End Slide ----
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1D, 0x1D, 0x1F)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Ready to build MAVIS 2.0"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Inter"

    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.7), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Mobile-first. Offline-capable. AI-powered.\nLet's replace the old Mavis — on your timeline."
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x86, 0x86, 0x8B)
    run2.font.name = "Inter"

    # Accent line
    line_shape = slide.shapes.add_shape(
        1, Inches(1.5), Inches(3.5), Inches(1.2), Pt(3)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0x0D, 0x94, 0x88)
    line_shape.line.fill.background()

    # Save
    prs.save(str(PPTX_PATH))
    print(f"\n  PPTX saved: {PPTX_PATH}")


if __name__ == "__main__":
    print("Taking screenshots...")
    take_screenshots()
    print("\nBuilding PPTX...")
    create_pptx()
    print("\nDone!")
